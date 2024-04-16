from pptx import Presentation
from pptx.util import Inches

SLD_LAYOUT_TITLE_AND_CONTENT = 1


def add_img_path(img_path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top)

    left = Inches(5)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save('test.pptx')


def main(old_ver, new_ver):
    prs = Presentation('2023_PPT_Template.pptx')

    # Opening page
    layout_open = prs.slide_masters[0].slide_layouts[1]
    slide = prs.slides.add_slide(layout_open)
    slide.shapes.title.text = 'Comparison Report of SMLib Versions'
    slide.placeholders[1].text = f'{old_ver} vs {new_ver}'

    # Disclaimer page
    layout_disclaimer = prs.slide_masters[1].slide_layouts[8]
    slide = prs.slides.add_slide(layout_disclaimer)

    # Comparison pages
    layout_comparison = prs.slide_masters[1].slide_layouts[3]  # left right two shape
    slide = prs.slides.add_slide(layout_comparison)
    slide.shapes.title.text = 'Comparison #1'
    #for shape in slide.placeholders:
    #    print('%d %s' % (shape.placeholder_format.idx, shape.name))
    slide.shapes.add_picture('0.png', Inches(0.5), Inches(1), width=Inches(6))
    slide.shapes.add_picture('0.png', Inches(6.7), Inches(1), width=Inches(6))
    x, y, cx, cy = Inches(2), Inches(4), Inches(8), Inches(1.5)
    table = slide.shapes.add_table(4, 4, x, y, cx, cy).table
    # write column headings
    table.cell(0, 1).text = old_ver
    table.cell(0, 2).text = new_ver
    table.cell(0, 3).text = 'Difference'
    # write body cells
    table.cell(1, 0).text = 'Srfm Free Edge'
    table.cell(2, 0).text = 'Srfm T-connect'
    table.cell(3, 0).text = 'Srfm Overlap'
    #slide.placeholders[2].insert_picture('0.png')

    # ---create presentation with 1 slide---
    #slide = prs.slides.add_slide(prs.slide_layouts[5])

    # ---add table to slide---
    # x, y, cx, cy = Inches(2), Inches(3), Inches(4), Inches(1.5)
    # shape = slide.shapes.add_table(3, 3, x, y, cx, cy)
    # shape_img = slide.shapes.add_picture('0.png', Inches(4), Inches(6))

    prs.save('test.pptx')


# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    main('SMLib 8.11', 'SMLib 8.12')

# See PyCharm help at https://www.jetbrains.com/help/pycharm/
